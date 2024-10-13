from flask import Flask, render_template, request, send_from_directory, send_file
from werkzeug.utils import secure_filename
import os
from pdf2docx import Converter
from PyPDF2 import PdfReader, PdfWriter
import fitz  # PyMuPDF
from PIL import Image
import zipfile
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import img2pdf

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'  # Folder to save uploaded files
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Sitemap route
@app.route('/sitemap.xml')
def sitemap():
    return send_from_directory(app.root_path, 'sitemap.xml')

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/blog')
def blog():
    return render_template('blog.html')

@app.route('/convert_pdf_to_word', methods=['POST'])
def convert_pdf_to_word():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Save the uploaded PDF file
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
    file.save(pdf_path)

    try:
        # Convert PDF to Word
        docx_path = pdf_path.replace('.pdf', '.docx')
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()

        # Send the converted file back
        return send_from_directory(app.config['UPLOAD_FOLDER'], os.path.basename(docx_path), as_attachment=True)
    except Exception as e:
        return f"Error converting PDF to Word: {str(e)}", 400

@app.route('/split_pdf', methods=['POST'])
def split_pdf():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Save the uploaded PDF file
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
    file.save(pdf_path)

    # Get the start and end page from the form
    start_page = int(request.form.get('start_page')) - 1  # Zero-indexed
    end_page = int(request.form.get('end_page')) - 1  # Zero-indexed

    try:
        # Read the PDF
        pdf_reader = PdfReader(pdf_path)

        # Ensure valid range
        if start_page < 0 or end_page >= len(pdf_reader.pages) or start_page > end_page:
            return "Invalid page range", 400

        # Create a zip file in memory
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for i in range(start_page, end_page + 1):
                pdf_writer = PdfWriter()
                pdf_writer.add_page(pdf_reader.pages[i])

                # Generate a filename for each split page
                split_page_filename = f'page_{i + 1}.pdf'

                # Write each split PDF to the zip file
                split_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], split_page_filename)
                with open(split_pdf_path, 'wb') as split_pdf_file:
                    pdf_writer.write(split_pdf_file)

                # Add the PDF file to the zip archive
                zip_file.write(split_pdf_path, arcname=split_page_filename)

        # Set the zip buffer's pointer to the start
        zip_buffer.seek(0)

        # Send the zip file as a downloadable file
        return send_file(zip_buffer, download_name='split_pages.zip', as_attachment=True)

    except Exception as e:
        return f"Error processing PDF: {str(e)}", 400

@app.route('/merge_pdf', methods=['POST'])
def merge_pdf():
    if 'files' not in request.files:
        return "No file part", 400
    files = request.files.getlist('files')
    pdf_writer = PdfWriter()

    try:
        # Save each uploaded PDF file and add to writer
        for file in files:
            if file.filename == '':
                return "No selected file", 400
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
            file.save(pdf_path)
            pdf_reader = PdfReader(pdf_path)
            for page in pdf_reader.pages:
                pdf_writer.add_page(page)

        merged_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], 'merged.pdf')
        with open(merged_pdf_path, 'wb') as out_file:
            pdf_writer.write(out_file)

        return send_from_directory(app.config['UPLOAD_FOLDER'], 'merged.pdf', as_attachment=True)

    except Exception as e:
        return f"Error merging PDFs: {str(e)}", 400

@app.route('/convert_pdf_to_image', methods=['POST'])
def convert_pdf_to_image():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Save the uploaded PDF file
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
    file.save(pdf_path)

    try:
        # Convert PDF to images using PyMuPDF
        image_paths = []
        pdf_document = fitz.open(pdf_path)

        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)
            pix = page.get_pixmap()
            img_path = os.path.join(app.config['UPLOAD_FOLDER'], f'page_{page_number + 1}.png')
            pix.save(img_path)
            image_paths.append(img_path)

        # Create a zip file in memory
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for img_path in image_paths:
                zip_file.write(img_path, arcname=os.path.basename(img_path))

        # Set the zip buffer's pointer to the start
        zip_buffer.seek(0)

        # Send the zip file as a downloadable file
        return send_file(zip_buffer, download_name='images.zip', as_attachment=True)

    except Exception as e:
        return f"Error converting PDF to images: {str(e)}", 400

@app.route('/convert_image_to_pdf', methods=['POST'])
def convert_image_to_pdf():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Save the uploaded image file
    image_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
    file.save(image_path)

    try:
        # Convert image to PDF using fitz
        pdf_path = image_path.replace('.jpg', '.pdf').replace('.png', '.pdf')
        pdf_document = fitz.open()  # Create a new PDF
        img = fitz.open(image_path)  # Open the image

        # Add the image to the PDF
        pdf_document.insert_page(-1)  # Add a new page
        pdf_document[-1].insert_image(pdf_document[-1].rect, filename=image_path)  # Insert image into the new page
        pdf_document.save(pdf_path)  # Save the PDF
        pdf_document.close()  # Close the PDF document

        return send_from_directory(app.config['UPLOAD_FOLDER'], os.path.basename(pdf_path), as_attachment=True)
    except Exception as e:
        return f"Error converting image to PDF: {str(e)}", 400

@app.route('/convert_pdf_to_ppt', methods=['POST'])
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Save the uploaded PDF file
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
    file.save(pdf_path)

    try:
        # Convert PDF to PPT
        ppt_path = pdf_path.replace('.pdf', '.pptx')
        presentation = Presentation()

        pdf_document = fitz.open(pdf_path)

        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)
            img = page.get_pixmap()  # Get a pixmap of the page
            img_path = os.path.join(app.config['UPLOAD_FOLDER'], f'page_{page_number + 1}.png')
            img.save(img_path)  # Save as PNG

            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            left = top = Inches(1)  # Adjust the image position as needed
            slide.shapes.add_picture(img_path, left, top, width=Inches(8))  # Add image to slide

        presentation.save(ppt_path)

        return send_from_directory(app.config['UPLOAD_FOLDER'], os.path.basename(ppt_path), as_attachment=True)
    except Exception as e:
        return f"Error converting PDF to PPT: {str(e)}", 400

@app.route('/convert_ppt_to_pdf', methods=['POST'])
def convert_ppt_to_pdf():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Save the uploaded PPT file
    ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(file.filename))
    file.save(ppt_path)

    try:
        # Convert PPT to PDF
        pdf_path = ppt_path.replace('.pptx', '.pdf')

        # Create a temporary folder to save slides as images
        temp_folder = os.path.join(app.config['UPLOAD_FOLDER'], 'temp_images')
        os.makedirs(temp_folder, exist_ok=True)

        # Load the presentation
        presentation = Presentation(ppt_path)

        # Save each slide as an image
        img_paths = []
        for slide_number, slide in enumerate(presentation.slides):
            img_path = os.path.join(temp_folder, f'slide_{slide_number + 1}.png')
            slide.shapes[0].element.getparent().getparent().make_image(img_path)  # Attempting to save as image

            # For now, we will use Pillow to create a blank image instead.
            img = Image.new('RGB', (800, 600), color='white')  # Create a blank image
            img.save(img_path)  # Save the blank image (replace with real image processing)

            img_paths.append(img_path)

        # Convert images to PDF using img2pdf
        with open(pdf_path, "wb") as f:
            f.write(img2pdf.convert(img_paths))

        # Clean up temporary images
        for img_file in img_paths:
            os.remove(img_file)
        os.rmdir(temp_folder)

        return send_from_directory(app.config['UPLOAD_FOLDER'], os.path.basename(pdf_path), as_attachment=True)
    except Exception as e:
        return f"Error converting PPT to PDF: {str(e)}", 400

@app.route('/static/<path:filename>')
def serve_static(filename):
    return send_from_directory('static', filename)

if __name__ == '__main__':
    app.run(debug=True)
